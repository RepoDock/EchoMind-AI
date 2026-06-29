from pptx import Presentation

from extractor.base_extractor import BaseExtractor


class PPTExtractor(BaseExtractor):

    def extract_text(self, file_path):

        presentation = Presentation(file_path)

        slides_text = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    if shape.text.strip():

                        slides_text.append(shape.text)

        text = "\n".join(slides_text)

        return {
            "text": text,
            "pages": len(presentation.slides),
            "title": "",
        }