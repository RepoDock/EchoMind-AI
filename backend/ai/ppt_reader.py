from pptx import Presentation
from pathlib import Path


def extract_ppt_text(file_path):

    presentation = Presentation(file_path)

    full_text = ""
    page_data = []

    for slide_no, slide in enumerate(presentation.slides, start=1):

        slide_text = ""

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"

        full_text += slide_text

        page_data.append({
            "page": slide_no,
            "text": slide_text
        })

    return {
        "text": full_text,
        "pages": len(presentation.slides),
        "title": Path(file_path).name,
        "page_data": page_data
    }